import re
import shutil
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
from pypdf import PdfReader

BASE_DIR = Path(__file__).resolve().parent.parent
DATA_DIR = BASE_DIR / "data"

MANUALS_DIR = DATA_DIR / "manuals"
GMAO_DIR = DATA_DIR / "gmao"
SCHEMAS_DIR = DATA_DIR / "schemas"
FAULT_CODES_DIR = DATA_DIR / "vims" if (DATA_DIR / "vims").exists() else DATA_DIR / "fault_codes"

# ── OCR optionnel ───────────────────────────────────────────────────────────
try:
    import pytesseract
    from pdf2image import convert_from_path
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("⚠️ OCR non disponible — pip install pytesseract pdf2image")

# ── DOCX optionnel ──────────────────────────────────────────────────────────
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️ python-docx non installé — pip install python-docx")

# ── PPTX optionnel ──────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx non installé — pip install python-pptx")

# ── ChromaDB + embeddings ───────────────────────────────────────────────────
try:
    import chromadb
    from chromadb.utils import embedding_functions
    CHROMA_AVAILABLE = True
except ImportError:
    CHROMA_AVAILABLE = False
    print("⚠️ chromadb non installé — fallback lexical activé")

CHROMA_PATH = DATA_DIR / "chroma_db"
EMBEDDING_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"


class RAGEngine:
    def __init__(self):
        self.documents: List[str] = []
        self.sources: List[str] = []
        self.metadatas: List[Dict[str, Any]] = []
        self.chroma_client = None
        self.collection = None
        self._init_tesseract()
        self._init_chroma()

    # ── Init ───────────────────────────────────────────────────────────────

    def _init_tesseract(self):
        if not OCR_AVAILABLE:
            return
        for p in [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        ]:
            if Path(p).exists():
                pytesseract.pytesseract.tesseract_cmd = p
                print(f"✅ Tesseract détecté : {p}")
                return

    def _init_chroma(self):
        if not CHROMA_AVAILABLE:
            return
        try:
            CHROMA_PATH.mkdir(parents=True, exist_ok=True)
            self.chroma_client = chromadb.PersistentClient(path=str(CHROMA_PATH))
            ef = embedding_functions.SentenceTransformerEmbeddingFunction(
                model_name=EMBEDDING_MODEL
            )
            self.collection = self.chroma_client.get_or_create_collection(
                name="mine_assist",
                embedding_function=ef,
                metadata={"hnsw:space": "cosine"},
            )
            print(f"✅ ChromaDB initialisé ({self.collection.count()} chunks existants)")
        except Exception as e:
            print(f"⚠️ Erreur ChromaDB : {e} — fallback lexical activé")
            self.chroma_client = None
            self.collection = None

    # ── Utilitaires texte ──────────────────────────────────────────────────

    def split_text(self, text: str, chunk_size: int = 700, overlap: int = 100) -> List[str]:
        text = (text or "").strip()
        if not text:
            return []

        sentences = re.compile(r'(?<=[.!?])\s+|\n{2,}').split(text)
        sentences = [s.strip() for s in sentences if s.strip()]

        chunks: List[str] = []
        current = ""
        for sentence in sentences:
            if len(current) + len(sentence) + 1 <= chunk_size:
                current = (current + " " + sentence).strip()
            else:
                if current:
                    chunks.append(current)
                overlap_text = current[-overlap:] if len(current) > overlap else current
                current = (overlap_text + " " + sentence).strip()

        if current:
            chunks.append(current)

        return chunks

    def _normalize_text(self, text: str) -> str:
        if not text:
            return ""
        s = str(text).lower()
        replacements = {
            "é": "e", "è": "e", "ê": "e", "ë": "e",
            "à": "a", "â": "a", "ä": "a",
            "ù": "u", "û": "u", "ü": "u",
            "ô": "o", "ö": "o",
            "î": "i", "ï": "i",
            "ç": "c",
        }
        for old, new in replacements.items():
            s = s.replace(old, new)
        s = re.sub(r"[^a-z0-9]+", " ", s)
        return re.sub(r"\s+", " ", s).strip()

    def _extract_page_markers(self, text: str) -> List[Tuple[int, str]]:
        """
        Découpe un texte PDF de type [Page X] ... en segments par page.
        Retourne [(page_num, texte_page), ...].
        """
        matches = list(re.finditer(r"\[Page\s+(\d+)\]", text))
        if not matches:
            return [(1, text)] if text.strip() else []

        parts: List[Tuple[int, str]] = []
        for i, match in enumerate(matches):
            page_num = int(match.group(1))
            start = match.end()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            page_text = text[start:end].strip()
            if page_text:
                parts.append((page_num, page_text))
        return parts

    # ── Loaders fichiers ────────────────────────────────────────────────────

    def load_pdf_text(self, path: str) -> List[Tuple[str, int]]:
        pages: List[Tuple[str, int]] = []
        try:
            reader = PdfReader(path)
            for i, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        pages.append((page_text.strip(), i + 1))
                except Exception as e:
                    print(f"Erreur lecture PDF {path}, page {i + 1}: {e}")
        except Exception as e:
            print(f"Erreur ouverture PDF {path}: {e}")
        return pages

    def load_pdf_ocr(self, path: str) -> List[Tuple[str, int]]:
        if not OCR_AVAILABLE:
            return []

        pages: List[Tuple[str, int]] = []
        temp_dir = tempfile.mkdtemp(prefix="mineassist_ocr_")
        try:
            images = convert_from_path(path, dpi=200, output_folder=temp_dir)
            for i, image in enumerate(images):
                try:
                    page_text = pytesseract.image_to_string(image, lang="eng+fra")
                    if page_text and page_text.strip():
                        pages.append((page_text.strip(), i + 1))
                except Exception as e:
                    print(f"Erreur OCR {path}, page {i + 1}: {e}")
        except Exception as e:
            print(f"Erreur OCR PDF {path}: {e}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)
        return pages

    def load_pdf(self, path: str) -> List[Tuple[str, int]]:
        pages = self.load_pdf_text(path)
        extracted_text_len = len(" ".join(t for t, _ in pages).strip())

        if extracted_text_len < 300:
            ocr_pages = self.load_pdf_ocr(path)
            ocr_text_len = len(" ".join(t for t, _ in ocr_pages).strip())
            if ocr_text_len > extracted_text_len:
                pages = ocr_pages
                print(f"ℹ️ OCR retenu pour {Path(path).name}")

        return pages

    def load_excel(self, path: str) -> str:
        """
        Charge un fichier Excel capteur.
        - Détecte automatiquement si l'en-tête est à la ligne 1 ou à la ligne 9 (format CAT export).
        - Pour les fichiers capteurs (colonnes Paramètres Diagnostic détectées),
          génère un résumé statistique sémantique par paramètre au lieu de dumper
          188 000 lignes brutes — ce qui noie le RAG et dilue le contexte.
        - Pour les autres xlsx (codes défauts, etc.), dump ligne par ligne comme avant.
        """
        filename = Path(path).name
        text = f"\n[Fichier: {filename}]\n"
        try:
            # ── Détection du format (en-tête ligne 1 ou ligne 9) ──────────────
            df_try1 = pd.read_excel(path, header=0, nrows=3)
            cols1 = [str(c).strip() for c in df_try1.columns]

            if "Paramètres Diagnostic" in cols1 or "Parametres Diagnostic" in cols1:
                df = pd.read_excel(path, header=0)
            else:
                # Format export CAT : 8 lignes de méta puis l'en-tête réel
                df = pd.read_excel(path, header=8)

            cols = [str(c).strip() for c in df.columns]
            df.columns = cols

            # ── Si c'est un fichier capteurs → résumé statistique ─────────────
            param_col = next((c for c in cols if "Param" in c and "Diagn" in c), None)
            avg_col   = next((c for c in cols if "moyenne" in c.lower() or "moy" in c.lower()), None)
            min_col   = next((c for c in cols if "minimale" in c.lower() or "min" in c.lower()), None)
            max_col   = next((c for c in cols if "maximale" in c.lower() or "max" in c.lower()), None)
            unit_col  = next((c for c in cols if "unit" in c.lower() or "unité" in c.lower()), None)
            time_col  = next((c for c in cols if "heure" in c.lower() or "date" in c.lower() or "horodatage" in c.lower()), None)

            if param_col and avg_col:
                df[avg_col] = pd.to_numeric(df[avg_col], errors="coerce")
                if min_col:
                    df[min_col] = pd.to_numeric(df[min_col], errors="coerce")
                if max_col:
                    df[max_col] = pd.to_numeric(df[max_col], errors="coerce")
                if time_col:
                    df[time_col] = pd.to_datetime(df[time_col], errors="coerce")

                # Période couverte
                if time_col and df[time_col].notna().any():
                    t_min = df[time_col].min().strftime("%Y-%m-%d")
                    t_max = df[time_col].max().strftime("%Y-%m-%d")
                    text += f"Période : {t_min} → {t_max}\n"

                text += f"Nombre de mesures : {len(df)}\n\n"
                text += "Résumé statistique par paramètre capteur :\n"

                grouped = df.groupby(param_col)
                for param_name, grp in grouped:
                    param_name = str(param_name).strip()
                    if not param_name or param_name.lower() == "nan":
                        continue
                    unite = ""
                    if unit_col:
                        u = grp[unit_col].dropna()
                        unite = str(u.iloc[0]).strip() if not u.empty else ""

                    vals = grp[avg_col].dropna()
                    if vals.empty:
                        continue

                    moy   = round(vals.mean(), 2)
                    vmin  = round(float(grp[min_col].min()), 2) if min_col else round(float(vals.min()), 2)
                    vmax  = round(float(grp[max_col].max()), 2) if max_col else round(float(vals.max()), 2)
                    nb    = len(vals)

                    # Résumé mensuel si la colonne temps existe
                    monthly = ""
                    if time_col and df[time_col].notna().any():
                        grp2 = grp.copy()
                        grp2["_mois"] = grp2[time_col].dt.to_period("M").astype(str)
                        m_summary = grp2.groupby("_mois")[avg_col].mean().round(1)
                        if len(m_summary) > 1:
                            monthly = " | Tendance mensuelle : " + ", ".join(
                                f"{m}={v}" for m, v in m_summary.items()
                            )

                    text += (
                        f"- {param_name} : moyenne={moy}{unite}, "
                        f"min={vmin}{unite}, max={vmax}{unite}, "
                        f"N={nb} mesures{monthly}\n"
                    )
                return text

            # ── Fichiers anomalies VIMS → format non-ambigu ──────────────────
            # Détecte si c'est un export anomalies (colonne "Code d'anomalie")
            anomalie_col = next(
                (c for c in cols if "anomalie" in c.lower() or "code d" in c.lower()), None
            )
            cid_col   = next((c for c in cols if c.lower().startswith("cid")), None)
            fmi_col   = next((c for c in cols if c.lower().startswith("fmi")), None)
            eid_col   = next((c for c in cols if c.lower().startswith("eid")), None)
            sev_col   = next((c for c in cols if "gravit" in c.lower()), None)
            occ_col   = next((c for c in cols if "occurrence" in c.lower()), None)
            date_col  = next((c for c in cols if "date" in c.lower()), None)
            src_col   = next((c for c in cols if "source" in c.lower()), None)

            if anomalie_col:
                # Résumé par code d'anomalie (compte occurrences, gravité max)
                text += "Liste des codes d'anomalie VIMS enregistrés :\n"
                text += "(ATTENTION : CID, FMI, EID sont des IDENTIFIANTS de code, pas des valeurs mesurées)\n\n"
                grouped = {}
                for _, row in df.iterrows():
                    code_nom = str(row.get(anomalie_col, "")).strip()
                    if not code_nom or code_nom.lower() == "nan":
                        continue
                    if code_nom not in grouped:
                        cid = row.get(cid_col, "") if cid_col else ""
                        fmi = row.get(fmi_col, "") if fmi_col else ""
                        eid = row.get(eid_col, "") if eid_col else ""
                        sev = row.get(sev_col, "") if sev_col else ""
                        src = row.get(src_col, "") if src_col else ""
                        # Formate les identifiants pour éviter toute confusion avec des valeurs
                        id_str = ""
                        if str(cid).strip() not in ["", "nan"]:
                            id_str += f"CID_id={int(float(cid))} FMI_id={int(float(fmi))}"
                        elif str(eid).strip() not in ["", "nan"]:
                            id_str += f"EID_id={int(float(eid))}"
                        grouped[code_nom] = {
                            "id": id_str,
                            "sev": str(sev).strip(),
                            "src": str(src).strip(),
                            "count": 0,
                        }
                    grouped[code_nom]["count"] += 1

                for nom, info in sorted(grouped.items(), key=lambda x: -x[1]["count"]):
                    text += (
                        f"[ANOMALIE] \"{nom}\" | {info['id']} | "
                        f"Occurrences={info['count']} | Gravité={info['sev']} | Source={info['src']}\n"
                    )
                return text

            # ── Autres xlsx sans structure reconnue → dump classique ─────────
            for _, row in df.iterrows():
                row_text = " | ".join(
                    f"{col}: {str(row[col])}" for col in df.columns if str(row[col]).strip()
                )
                if row_text.strip():
                    text += row_text + "\n"

        except Exception as e:
            print(f"Erreur lecture Excel {path}: {e}")
        return text

    def load_csv(self, path: str) -> str:
        try:
            df = pd.read_csv(path).fillna("").astype(str)
            return df.to_string(index=False)
        except Exception as e:
            print(f"Erreur lecture CSV {path}: {e}")
            return ""

    def load_docx(self, path: str) -> List[Tuple[str, int]]:
        if not DOCX_AVAILABLE:
            return []

        sections: List[Tuple[str, int]] = []
        try:
            doc = DocxDocument(str(path))
            current_section: List[str] = []
            section_num = 1

            for para in doc.paragraphs:
                txt = para.text.strip()
                if not txt:
                    continue

                style_name = para.style.name if para.style and para.style.name else ""
                is_heading = (
                    style_name.startswith("Heading")
                    or style_name.startswith("Title")
                    or (len(txt) < 120 and txt.isupper())
                )

                if is_heading and current_section:
                    full_text = "\n".join(current_section).strip()
                    if full_text:
                        sections.append((full_text, section_num))
                    section_num += 1
                    current_section = [txt]
                else:
                    current_section.append(txt)

            if current_section:
                full_text = "\n".join(current_section).strip()
                if full_text:
                    sections.append((full_text, section_num))

            for i, table in enumerate(doc.tables):
                table_lines: List[str] = []
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells:
                        table_lines.append(" | ".join(row_cells))
                if table_lines:
                    sections.append(("\n".join(table_lines), 1000 + i))

        except Exception as e:
            print(f"Erreur lecture DOCX {path}: {e}")

        return sections

    def load_pptx(self, path: str) -> List[Tuple[str, int]]:
        if not PPTX_AVAILABLE:
            return []

        slides: List[Tuple[str, int]] = []
        try:
            prs = Presentation(path)
            for i, slide in enumerate(prs.slides):
                parts: List[str] = []

                if slide.shapes.title and slide.shapes.title.text.strip():
                    parts.append(f"[Titre] {slide.shapes.title.text.strip()}")

                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        title_txt = slide.shapes.title.text.strip() if slide.shapes.title else ""
                        if txt and txt != title_txt:
                            parts.append(txt)

                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes] {notes}")

                full_text = "\n".join(parts).strip()
                if full_text:
                    slides.append((full_text, i + 1))

        except Exception as e:
            print(f"Erreur lecture PPTX {path}: {e}")

        return slides

    # ── Indexation ─────────────────────────────────────────────────────────

    def _add_chunks(self, chunks_with_pages: List[Tuple[str, int]], filename: str, doc_type: str):
        for text, page in chunks_with_pages:
            sub_chunks = self.split_text(text)
            for chunk in sub_chunks:
                self.documents.append(chunk)
                self.sources.append(filename)
                self.metadatas.append({
                    "source": filename,
                    "page": page,
                    "type": doc_type,
                })

    def load_documents_from_folder(self, folder_path: Path, allowed_extensions: List[str]):
        if not folder_path.exists():
            print(f"⚠️ Dossier introuvable : {folder_path}")
            return

        # Scan récursif : inclut les sous-dossiers (gmao/capteurs/, gmao/anomalies/, etc.)
        for file_path in sorted(folder_path.rglob("*")):
            if not file_path.is_file():
                continue
            ext = file_path.suffix.lower()
            if ext not in allowed_extensions:
                continue

            try:
                name = file_path.name

                if ext == ".pdf":
                    pages = self.load_pdf(str(file_path))
                    self._add_chunks(pages, name, "pdf")
                    print(f"📄 PDF : {name} ({len(pages)} pages)")

                elif ext == ".docx":
                    sections = self.load_docx(str(file_path))
                    self._add_chunks(sections, name, "docx")
                    print(f"📝 DOCX : {name} ({len(sections)} sections)")

                elif ext in [".pptx", ".ppt"]:
                    slides = self.load_pptx(str(file_path))
                    self._add_chunks(slides, name, "pptx")
                    print(f"📊 PPTX : {name} ({len(slides)} diapositives)")

                elif ext in [".xlsx", ".xls"]:
                    text = self.load_excel(str(file_path))
                    chunks = self.split_text(text)
                    for chunk in chunks:
                        self.documents.append(chunk)
                        self.sources.append(name)
                        self.metadatas.append({
                            "source": name,
                            "page": 0,
                            "type": "excel",
                        })
                    print(f"📈 Excel : {name} ({len(chunks)} chunks)")

                elif ext == ".csv":
                    text = self.load_csv(str(file_path))
                    chunks = self.split_text(text)
                    for chunk in chunks:
                        self.documents.append(chunk)
                        self.sources.append(name)
                        self.metadatas.append({
                            "source": name,
                            "page": 0,
                            "type": "csv",
                        })
                    print(f"📄 CSV : {name} ({len(chunks)} chunks)")

            except Exception as e:
                print(f"Erreur lecture fichier {file_path.name}: {e}")

    def load_documents(self):
        self.documents = []
        self.sources = []
        self.metadatas = []

        self.load_documents_from_folder(MANUALS_DIR, [".pdf", ".docx", ".pptx", ".ppt"])
        self.load_documents_from_folder(GMAO_DIR, [".xlsx", ".xls", ".csv"])
        self.load_documents_from_folder(SCHEMAS_DIR, [".pdf", ".docx", ".pptx", ".ppt"])
        self.load_documents_from_folder(FAULT_CODES_DIR, [".pdf", ".docx", ".xlsx", ".xls", ".csv"])

        print(f"📚 Total : {len(self.documents)} chunks de {len(set(self.sources))} sources")

    def index_all(self):
        self.load_documents()

        if self.collection is not None and self.documents:
            try:
                existing_ids = set(self.collection.get()["ids"])
            except Exception:
                existing_ids = set()

            new_docs: List[str] = []
            new_ids: List[str] = []
            new_metas: List[Dict[str, Any]] = []

            for i, (doc, meta) in enumerate(zip(self.documents, self.metadatas)):
                doc_id = f"{meta['source']}__p{meta['page']}__chunk_{i}"
                if doc_id not in existing_ids:
                    new_docs.append(doc)
                    new_ids.append(doc_id)
                    new_metas.append(meta)

            if new_docs:
                batch_size = 100
                for start in range(0, len(new_docs), batch_size):
                    self.collection.add(
                        documents=new_docs[start:start + batch_size],
                        ids=new_ids[start:start + batch_size],
                        metadatas=new_metas[start:start + batch_size],
                    )
                print(f"✅ {len(new_docs)} nouveaux chunks indexés dans ChromaDB")
            else:
                print("ℹ️ Aucun nouveau document à indexer")

        return {
            "documents_indexed": len(self.documents),
            "sources_indexed": len(set(self.sources)),
        }

    # ── Recherche ──────────────────────────────────────────────────────────

    def _format_source_label(self, meta: Dict[str, Any]) -> str:
        src = meta.get("source", "")
        page = meta.get("page", 0)
        typ = meta.get("type", "")

        if not page or page == 0:
            return src
        if typ == "pptx":
            return f"{src} (diapo {page})"
        if typ == "docx":
            return f"{src} (section {page})"
        return f"{src} (page {page})"

    def _technical_translations(self) -> Dict[str, str]:
        return {
            "hydraulique": "hydraulic",
            "moteur": "engine",
            "frein": "brake",
            "huile": "oil",
            "pression": "pressure",
            "pompe": "pump",
            "filtre": "filter",
            "temperature": "temperature",
            "temp": "temperature",
            "defaut": "fault",
            "diagnostic": "diagnostic",
            "capteur": "sensor",
            "vanne": "valve",
            "transmission": "transmission",
            "embrayage": "clutch",
            "essieu": "axle",
            "refroidissement": "cooling",
            "turbo": "turbo",
        }

    # Mots-clés qui signalent une question sur des valeurs capteurs terrain
    _SENSOR_KEYWORDS = [
        "température", "temperature", "pression", "pressure", "régime", "regime",
        "vitesse", "tension", "courant", "débit", "niveau", "normale", "normal",
        "valeur", "mesure", "capteur", "moyenne", "min", "max", "seuil",
        "liquide refroidissement", "huile moteur", "huile direction", "huile freinage",
        "echappement", "échappement", "convertisseur", "impeller", "hydraulique",
        "essieux", "pto", "regime moteur", "régime moteur",
    ]

    def _is_sensor_query(self, query: str) -> bool:
        q = query.lower()
        return any(kw in q for kw in self._SENSOR_KEYWORDS)

    def _get_excel_capteur_chunks_from_chroma(self, query: str) -> List[Tuple[str, Dict]]:
        """
        Interroge ChromaDB directement avec where={"type": "excel"} pour récupérer
        les chunks capteurs pertinents. Fiable même quand self.documents est vide
        (cas normal au démarrage FastAPI sans appel à index_all).
        """
        if self.collection is None or self.collection.count() == 0:
            return []
        try:
            # Compte d'abord les chunks Excel disponibles
            all_excel = self.collection.get(where={"type": "excel"})
            if not all_excel or not all_excel.get("ids"):
                return []
            n_excel = len(all_excel["ids"])
            if n_excel == 0:
                return []

            # Recherche sémantique filtrée sur type=excel uniquement
            results = self.collection.query(
                query_texts=[query],
                n_results=min(n_excel, 6),
                where={"type": "excel"},
                include=["documents", "metadatas", "distances"],
            )
            docs   = results["documents"][0]
            metas  = results["metadatas"][0]
            dists  = results["distances"][0]

            # Pas de seuil de distance strict : on prend tout ce qui contient des mots clés
            query_norm  = self._normalize_text(query)
            query_words = [w for w in query_norm.split() if len(w) > 3]
            matched = []
            for doc, meta, dist in zip(docs, metas, dists):
                doc_norm = self._normalize_text(doc)
                kw_score = sum(1 for w in query_words if w in doc_norm)
                if kw_score > 0:
                    matched.append((kw_score, doc, meta))

            matched.sort(key=lambda x: x[0], reverse=True)
            return [(doc, meta) for _, doc, meta in matched[:3]]

        except Exception as e:
            print(f"⚠️ Erreur requête Excel ChromaDB : {e}")
            return []

    def build_context(self, query: str, top_k: int = 5, max_chars: int = 6000) -> Tuple[str, List[str]]:
        if not query:
            return "", []

        # ── Injection garantie des chunks Excel capteurs pour questions terrain ─
        excel_chunks: List[Tuple[str, Dict]] = []
        if self._is_sensor_query(query):
            excel_chunks = self._get_excel_capteur_chunks_from_chroma(query)
            if excel_chunks:
                print(f"📊 Injection Excel capteurs : {len(excel_chunks)} chunk(s) injectés en tête de contexte")
            else:
                print("⚠️ Aucun chunk Excel capteur trouvé dans ChromaDB — vérifier l'indexation")

        # ── Recherche sémantique via ChromaDB (PDF + autres) ──────────────
        if self.collection is not None and self.collection.count() > 0:
            try:
                results = self.collection.query(
                    query_texts=[query],
                    n_results=min(top_k, self.collection.count()),
                    include=["documents", "metadatas", "distances"],
                )

                docs      = results["documents"][0]
                metas     = results["metadatas"][0]
                distances = results["distances"][0]

                # Exclure les chunks Excel déjà injectés ci-dessus
                semantic = [
                    (doc, meta)
                    for doc, meta, dist in zip(docs, metas, distances)
                    if dist < 0.78 and meta.get("type") != "excel"
                ]

                # Excel en tête → contexte métier précis avant le contexte documentaire
                combined = excel_chunks + semantic

                if combined:
                    context = "\n\n---\n\n".join(doc for doc, _ in combined)[:max_chars]
                    sources: List[str] = []
                    seen = set()
                    for _, meta in combined:
                        label = self._format_source_label(meta)
                        if label not in seen:
                            seen.add(label)
                            sources.append(label)
                    return context, sources

                print("ℹ️ Aucun chunk pertinent trouvé (distance > 0.78)")

            except Exception as e:
                print(f"⚠️ Erreur ChromaDB query : {e} — fallback lexical")

        # ── Fallback lexical ───────────────────────────────────────────────
        if not self.documents:
            return "", []

        query_norm = self._normalize_text(query)
        query_words = query_norm.split()
        translations = self._technical_translations()
        scored_docs: List[Tuple[int, str, Dict[str, Any]]] = []

        for i, doc in enumerate(self.documents):
            score = 0
            doc_norm = self._normalize_text(doc)

            for word in query_words:
                if not word:
                    continue
                if word in doc_norm:
                    score += 2
                if len(word) > 4 and word.rstrip("e") in doc_norm:
                    score += 1
                en_word = translations.get(word)
                if en_word and en_word in doc_norm:
                    score += 3

            # bonus pour codes / termes fréquents
            if any(code in doc for code in re.findall(r"\b(?:MID|CID|FMI|ECM|VIMS)\s*\d+\b", query.upper())):
                score += 4

            if score > 0:
                scored_docs.append((score, doc, self.metadatas[i]))

        scored_docs.sort(key=lambda x: x[0], reverse=True)
        selected = scored_docs[:top_k]
        if not selected:
            return "", []

        context = "\n\n---\n\n".join(doc for _, doc, _ in selected)[:max_chars]
        sources: List[str] = []
        seen = set()
        for _, _, meta in selected:
            label = self._format_source_label(meta)
            if label not in seen:
                seen.add(label)
                sources.append(label)

        return context, sources
